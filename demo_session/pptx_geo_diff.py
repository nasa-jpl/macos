"""Diff shape GEOMETRY (position/size) and FONT STYLE between two pptx decks.

Usage: python3 pptx_geo_diff.py baseline.pptx edited.pptx

Companion to pptx_text_diff.py (which sees only text).  Shapes are
matched within each slide by a normalized text prefix (text frames) or
by picture blob hash (images), falling back to shape order.  Reports
moves/resizes in inches (suppressed below 0.02") plus font size/name
changes per shape.  STANDING RULE (Dave 2026-08-31): run this AND the
text diff against the current baseline, reading the FULL output (never
piped through head), before every edit-deck re-sync.
"""
import sys, hashlib
from pptx import Presentation
from pptx.util import Emu

EMU_IN = 914400.0
TOL = 0.02  # inches


def style_of(sh):
    if not sh.has_text_frame:
        return None
    sizes, names = set(), set()
    for p in sh.text_frame.paragraphs:
        for r in p.runs:
            if r.font.size is not None:
                sizes.add(round(r.font.size.pt, 1))
            if r.font.name:
                names.add(r.font.name)
    return (tuple(sorted(sizes)), tuple(sorted(names)))


def key_of(sh, seq):
    if sh.shape_type == 13:  # PICTURE
        try:
            h = hashlib.md5(sh.image.blob).hexdigest()[:8]
            return f"pic:{h}"
        except Exception:
            return f"pic#{seq}"
    if sh.has_text_frame:
        t = " ".join(sh.text_frame.text.split())[:48]
        if t:
            return f"txt:{t}"
    return f"{sh.shape_type}#{seq}"


def shapes_of(path):
    prs = Presentation(path)
    out = []
    for sl in prs.slides:
        d = {}
        for i, sh in enumerate(sl.shapes):
            k = key_of(sh, i)
            n = 2
            base = k
            while k in d:  # disambiguate duplicates in slide order
                k = f"{base}({n})"; n += 1
            d[k] = (sh.left, sh.top, sh.width, sh.height, style_of(sh))
        out.append(d)
    return out


def fmt(emu):
    return f"{(emu or 0)/EMU_IN:.2f}"


a = shapes_of(sys.argv[1]); b = shapes_of(sys.argv[2])
for si in range(max(len(a), len(b))):
    A = a[si] if si < len(a) else {}
    B = b[si] if si < len(b) else {}
    lines = []
    for k in A:
        if k not in B:
            lines.append(f"  only in OLD: {k}")
            continue
        (l0, t0, w0, h0, s0), (l1, t1, w1, h1, s1) = A[k], B[k]
        dl, dt = (l1 - l0) / EMU_IN, (t1 - t0) / EMU_IN
        dw, dh = (w1 - w0) / EMU_IN, (h1 - h0) / EMU_IN
        parts = []
        if abs(dl) > TOL or abs(dt) > TOL:
            parts.append(f"moved ({dl:+.2f},{dt:+.2f})\" to ({fmt(l1)},{fmt(t1)})")
        if abs(dw) > TOL or abs(dh) > TOL:
            parts.append(f"resized ({dw:+.2f},{dh:+.2f})\" to {fmt(w1)}x{fmt(h1)}")
        if s0 != s1:
            parts.append(f"font {s0} -> {s1}")
        if parts:
            lines.append(f"  {k}: " + "; ".join(parts))
    for k in B:
        if k not in A:
            lines.append(f"  only in NEW: {k} at ({fmt(B[k][0])},{fmt(B[k][1])}) {fmt(B[k][2])}x{fmt(B[k][3])}")
    if lines:
        print(f"\n=== SLIDE {si+1}:")
        print("\n".join(lines))
