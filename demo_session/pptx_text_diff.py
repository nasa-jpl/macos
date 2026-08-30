import sys, difflib
from pptx import Presentation

def texts(path):
    prs = Presentation(path)
    out = []
    for sl in prs.slides:
        st = []
        for sh in sl.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                st.append(sh.text_frame.text)
        out.append(st)
    return out

a = texts(sys.argv[1]); b = texts(sys.argv[2])
for si in range(max(len(a), len(b))):
    A = a[si] if si < len(a) else []
    B = b[si] if si < len(b) else []
    for i in range(max(len(A), len(B))):
        ta = A[i] if i < len(A) else ""
        tb = B[i] if i < len(B) else ""
        if ta != tb:
            print(f"\n=== SLIDE {si+1} shape {i}:")
            for l in difflib.unified_diff(ta.split("\n"), tb.split("\n"),
                                          lineterm="", n=0):
                if l[:3] in ("---", "+++"): continue
                print(" " + l)
