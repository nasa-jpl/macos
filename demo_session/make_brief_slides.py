#!/usr/bin/env python3
"""Build a .pptx deck from a slide-markdown source.

Usage:  python3 make_brief_slides.py deck_fang.md [out.pptx]

Adapted from mmacos/design/examples/e2e/make_e2e_report.py with two
additions: markdown pipe-tables render as native pptx tables, and
**bold** / `code` markers are stripped from body text.  Conventions:
  # Title            -> title slide (following bare lines: big/paras; ~ notes)
  ## Head | subhead  -> new slide
  ::: left|right|mid|full   column select
  ::: stack          -> boxed flow items ("- head :: sub")
  ::: cards N        -> card grid ("- line1 :: line2 :: ...")
  - item             -> bullet
  ~ text             -> gray footnote
  ![cap](figs/x.png){h=2.5}  -> image, optional height cap (inches)
  | a | b |          -> table (header row + body; --- separator dropped)
Needs python-pptx + Pillow.
"""
import json
import os
import re
import sys
from math import ceil

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

MD = os.path.abspath(sys.argv[1])
HERE = os.path.dirname(MD)
OUT = os.path.abspath(sys.argv[2]) if len(sys.argv) > 2 else \
    os.path.splitext(MD)[0] + ".pptx"

# Geometry-override sidecar (<md stem>.geo.json): carries manual layout
# edits (Dave's pptx repositioning passes) across rebuilds, since the md
# holds no per-shape geometry.  Structure:
#   { "<slide title>": { "<block key>": {"x":in,"y":in,"w":in,"h":in,
#                                        "cx":in,"cy":in,"cw":in}, ... } }
# Block keys: "img:<fig basename>", "txt:<leading text of the block>",
# "code:N" / "table:N" (Nth code/table block on the slide, 1-based).
# JSON keys match as PREFIXES of the computed block key.  x/y/w place the
# block (y pins the flow cursor); w+h on an img place the picture EXACTLY
# (no fit-and-center); cx/cy/cw place its caption independently.
# Recover new edits with pptx_geo_diff.py (text edits: pptx_text_diff.py).
GEO_PATH = os.path.splitext(MD)[0] + ".geo.json"
GEO = {}
if os.path.exists(GEO_PATH):
    with open(GEO_PATH, encoding="utf-8") as f:
        GEO = json.load(f)

INK = RGBColor(0x28, 0x28, 0x28)
ACCENT = RGBColor(0x1F, 0x4E, 0x79)
GRAY = RGBColor(0x6E, 0x6E, 0x6E)
LIGHT = RGBColor(0xEC, 0xF1, 0xF6)
RULE = RGBColor(0xB8, 0xC6, 0xD4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SW_IN, SH_IN = 13.333, 7.5
BOTTOM = 7.32


def clean(s):
    s = re.sub(r"\*\*(.+?)\*\*", r"\1", s)
    s = re.sub(r"\*(.+?)\*", r"\1", s)
    s = s.replace("`", "")
    return s


# --------------------------------------------------------------- parsing
def parse(md_text):
    md_text = re.sub(r"<!--.*?-->", "", md_text, flags=re.S)
    lines = md_text.split("\n")
    title_slide = {"title": None, "big": [], "paras": [], "notes": []}
    slides = []
    cur = None
    col = "full"
    mode = None
    i = 0
    while i < len(lines):
        ln = lines[i]
        s = ln.strip()
        if s.startswith("## "):
            head = s[3:]
            t, _, sub = head.partition("|")
            cur = {"title": clean(t.strip()), "sub": clean(sub.strip()),
                   "blocks": []}
            slides.append(cur)
            col, mode = "full", None
        elif s.startswith("# ") and cur is None:
            title_slide["title"] = clean(s[2:].strip())
        elif s.startswith(":::"):
            tok = s[3:].strip().split()
            if tok and tok[0] in ("left", "right", "mid", "full"):
                col, mode = tok[0], None
            elif tok and tok[0] == "stack":
                mode = "stack"
                cur["blocks"].append((col, "stack", []))
            elif tok and tok[0] == "cards":
                mode = ("cards", int(tok[1]) if len(tok) > 1 else 3)
                cur["blocks"].append((col, "cards", (mode[1], [])))
        elif s.startswith("|"):
            rows = []
            while i < len(lines) and lines[i].strip().startswith("|"):
                cells = [clean(c.strip())
                         for c in lines[i].strip().strip("|").split("|")]
                if not all(re.fullmatch(r":?-{2,}:?", c) for c in cells):
                    rows.append(cells)
                i += 1
            i -= 1
            cur["blocks"].append((col, "table", rows))
            mode = None
        elif s.startswith("```"):
            code = []
            i += 1
            while i < len(lines) and not lines[i].strip().startswith("```"):
                code.append(lines[i].rstrip())
                i += 1
            cur["blocks"].append((col, "code", code))
        elif s.startswith("### "):
            cur["blocks"].append((col, "h3", clean(s[4:].strip())))
            mode = None
        elif s.startswith("-- "):
            # sub-bullet: indented, smaller, gray, under the current bullets
            item = "\x01" + clean(s[3:].strip())
            if cur["blocks"] and cur["blocks"][-1][0] == col \
                    and cur["blocks"][-1][1] == "bullets":
                cur["blocks"][-1][2].append(item)
            else:
                cur["blocks"].append((col, "bullets", [item]))
        elif s.startswith("- "):
            item = clean(s[2:].strip())
            if mode == "stack":
                head, _, sub = item.partition("::")
                cur["blocks"][-1][2].append((head.strip(), sub.strip()))
            elif isinstance(mode, tuple) and mode[0] == "cards":
                cur["blocks"][-1][2][1].append(
                    [p.strip() for p in item.split("::")])
            else:
                if cur["blocks"] and cur["blocks"][-1][0] == col \
                        and cur["blocks"][-1][1] == "bullets":
                    cur["blocks"][-1][2].append(item)
                else:
                    cur["blocks"].append((col, "bullets", [item]))
        elif s.startswith("~ "):
            if cur is None:
                title_slide["notes"].append(clean(s[2:].strip()))
            else:
                cur["blocks"].append((col, "note", clean(s[2:].strip())))
        elif s.startswith("!["):
            m = re.match(r"!\[(.*)\]\((.*?)\)(?:\{h=([\d.]+)\})?", s)
            if m:
                cap, path, h = clean(m.group(1)), m.group(2), m.group(3)
                cur["blocks"].append(
                    (col, "img", (cap, path, float(h) if h else None)))
        elif s:
            if cur is None:
                if not title_slide["big"]:
                    title_slide["big"].append(clean(s))
                else:
                    title_slide["paras"].append(clean(s))
            else:
                cur["blocks"].append((col, "para", clean(s)))
                mode = None
        i += 1
    return title_slide, slides


# ------------------------------------------------------------- rendering
prs = Presentation()
prs.slide_width, prs.slide_height = Inches(SW_IN), Inches(SH_IN)
BLANK = prs.slide_layouts[6]


def tb(sl, x, y, w, h):
    b = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = b.text_frame
    tf.word_wrap = True
    return tf


def para(tf, text, size=12, bold=False, color=INK, first=False, mono=False,
         space_after=4, align=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.text = text
    p.space_after = Pt(space_after)
    if align:
        p.alignment = align
    for r in p.runs:
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = "Consolas" if mono else "Arial"
    return p


def est_text_h(text, w_in, size):
    cpl = max(10, w_in * 72.0 / (0.5 * size))
    n = max(1, ceil(len(text) / cpl))
    return n * size * 1.30 / 72.0 + 0.07


def body_size(w_in):
    # Full-width bullet blocks read at 15 pt (Dave's 2026-08-30 edit pass:
    # 12 pt full-width reads as report text at room scale); columns keep
    # the tighter sizes.
    if w_in >= 10:
        return 15
    return 12 if w_in >= 5.5 else (11.5 if w_in >= 4.5 else 11)


def add_title(sl, text, sub):
    tf = tb(sl, 0.45, 0.22, SW_IN - 0.9, 1.0)
    para(tf, text, size=24, bold=True, color=ACCENT, first=True)
    if sub:
        para(tf, sub, size=12.5, color=GRAY, space_after=0)
    y = 1.24 if sub else 1.0
    ln = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.45), Inches(y),
                             Inches(SW_IN - 0.9), Pt(1.6))
    ln.fill.solid(); ln.fill.fore_color.rgb = RULE
    ln.line.fill.background(); ln.shadow.inherit = False
    return y + 0.18


def place_img(sl, path, x_in, y_in, w_in, h_in, caption, exact=False,
              cap_ov=None):
    full = os.path.join(HERE, path)
    if exact:
        # sidecar-pinned geometry: draw at (x,y) w x h verbatim
        pic = sl.shapes.add_picture(full, Inches(x_in), Inches(y_in),
                                    Inches(w_in), Inches(h_in))
        if cap_ov:   # crop fractions (Impress/PowerPoint semantics)
            for attr, key in (("crop_left", "cl"), ("crop_right", "cr"),
                              ("crop_top", "ct"), ("crop_bottom", "cb")):
                if key in cap_ov:
                    setattr(pic, attr, cap_ov[key])
        yy = y_in + h_in
        cx, cw = x_in, w_in
    else:
        iw, ih = Image.open(full).size
        cap_h0 = est_text_h(caption, w_in, 10) if caption else 0.0
        box_h = max(0.4, h_in - cap_h0)
        # fit box (w_in x box_h) preserving aspect
        sc = min(w_in / iw, box_h / ih)
        w, h = iw * sc, ih * sc
        px = x_in + (w_in - w) / 2.0
        sl.shapes.add_picture(full, Inches(px), Inches(y_in),
                              Inches(w), Inches(h))
        yy = y_in + h
        cx, cw = x_in, w_in
    if caption:
        cy = yy + 0.02
        cfs = 10
        if cap_ov:
            cx = cap_ov.get("cx", cx)
            cy = cap_ov.get("cy", cy)
            cw = cap_ov.get("cw", cw)
            cfs = cap_ov.get("cfs", cfs)
        cap_h = est_text_h(caption, cw, cfs)
        tf = tb(sl, cx, cy, cw, cap_h)
        para(tf, caption, size=cfs, color=GRAY, first=True,
             align=PP_ALIGN.CENTER, space_after=0)
        yy = cy + cap_h + 0.04
    return yy


def place_table(sl, rows, x, y, w):
    nr = len(rows)
    nc = max(len(r) for r in rows)
    size = 11 if w >= 6 else 10
    maxlens = [max((len(r[c]) if c < len(r) else 0) for r in rows) or 1
               for c in range(nc)]
    tot = float(sum(maxlens))
    cw = [w * (m / tot) for m in maxlens]
    # Naturally-narrow columns get their full unwrapped width when there is
    # room; long-text columns wrap anyway, so they absorb the difference.
    # (Otherwise a short label column beside paragraph cells is squeezed
    # into wrapping -- the proportional share is the fallback.)
    need = [m * 0.52 * size / 72.0 + 0.16 for m in maxlens]
    small = [c for c in range(nc) if need[c] <= 0.22 * w]
    if small and len(small) < nc and sum(need[c] for c in small) < 0.5 * w:
        rem = w - sum(need[c] for c in small)
        tot2 = float(sum(maxlens[c] for c in range(nc)
                         if c not in small)) or 1.0
        cw = [need[c] if c in small else rem * maxlens[c] / tot2
              for c in range(nc)]
    # per-row height from the worst wrapped cell (cpl mirrors est_text_h)
    rhs = []
    for r in rows:
        lines = 1
        for c in range(nc):
            txt = r[c] if c < len(r) else ""
            cpl = max(8, cw[c] * 72.0 / (0.52 * size))
            lines = max(lines, ceil(len(txt) / cpl))
        rhs.append(lines * size * 1.35 / 72.0 + 0.10)
    shape = sl.shapes.add_table(nr, nc, Inches(x), Inches(y),
                                Inches(w), Inches(sum(rhs)))
    tbl = shape.table
    for ir in range(nr):
        tbl.rows[ir].height = Emu(int(Inches(rhs[ir])))
    for c in range(nc):
        tbl.columns[c].width = Emu(int(Inches(cw[c])))
    for ir, r in enumerate(rows):
        for c in range(nc):
            cell = tbl.cell(ir, c)
            cell.margin_left = cell.margin_right = Inches(0.06)
            cell.margin_top = cell.margin_bottom = Inches(0.02)
            tf_ = cell.text_frame
            tf_.word_wrap = True
            p = tf_.paragraphs[0]
            p.text = r[c] if c < len(r) else ""
            for run in p.runs:
                run.font.size = Pt(size)
                run.font.name = "Arial"
                run.font.bold = (ir == 0)
                run.font.color.rgb = WHITE if ir == 0 else INK
            cell.fill.solid()
            cell.fill.fore_color.rgb = ACCENT if ir == 0 else \
                (WHITE if ir % 2 else LIGHT)
    return y + sum(rhs) + 0.16


def block_key(typ, payload, counters):
    if typ == "img":
        return "img:" + os.path.basename(payload[1])
    if typ in ("code", "table", "stack", "cards"):
        counters[typ] = counters.get(typ, 0) + 1
        return "%s:%d" % (typ, counters[typ])
    if typ == "bullets":
        return "txt:" + payload[0].lstrip("\x01")
    if typ in ("para", "note", "h3"):
        return "txt:" + str(payload)
    return typ


def render_slide(spec):
    sl = prs.slides.add_slide(BLANK)
    top = add_title(sl, spec["title"], spec["sub"])
    three = any(b[0] == "mid" for b in spec["blocks"])
    if three:
        cols = {"left": (0.45, 4.15), "mid": (4.72, 4.1),
                "right": (8.94, 3.95), "full": (0.45, SW_IN - 0.9)}
    else:
        cols = {"left": (0.45, 6.35), "right": (7.0, 5.9),
                "full": (0.45, SW_IN - 0.9)}
    cur = {c: top for c in cols}
    sov = GEO.get(spec["title"], {})
    counters = {}

    def sync_full():
        y = max(cur.values())
        for c in cur:
            cur[c] = y
        return y

    for colname, typ, payload in spec["blocks"]:
        if colname == "full":
            sync_full()
        x, w = cols[colname]
        y = cur[colname]
        bk = block_key(typ, payload, counters)
        ov = next((v for k, v in sov.items() if bk.startswith(k)), None)
        if ov:
            x = ov.get("x", x)
            w = ov.get("w", w)
            y = ov.get("y", y)
        if typ == "h3":
            tf = tb(sl, x, y, w, 0.35)
            para(tf, payload, size=13.5, bold=True, color=ACCENT,
                 first=True, space_after=0)
            y += 0.34
        elif typ in ("bullets", "para", "note"):
            items = payload if typ == "bullets" else [payload]
            size = body_size(w) if typ != "note" else 10.5
            if ov and "fs" in ov:   # sidecar font-size override (points)
                size = ov["fs"]
            color = INK if typ != "note" else GRAY
            tf = tb(sl, x, y, w, 0.3)
            for j, it in enumerate(items):
                if typ == "bullets" and it.startswith("\x01"):
                    ss = size - (2 if size >= 15 else 1.5)
                    txt = "–  " + it[1:]
                    p = para(tf, txt, size=ss, color=GRAY,
                             first=(j == 0), space_after=5)
                    # hanging indent so wrapped lines align after the dash
                    pPr = p._p.get_or_add_pPr()
                    pPr.set("marL", "594360")   # 0.65" left margin
                    pPr.set("indent", "-137160")  # dash hangs 0.15"
                    y += est_text_h(txt, w - 0.6, ss) + 5 / 72.0
                    continue
                txt = ("•  " + it) if typ == "bullets" else it
                para(tf, txt, size=size, color=color, first=(j == 0),
                     space_after=6)
                y += est_text_h(txt, w, size) + 6 / 72.0
            y += 0.04
        elif typ == "table":
            y = place_table(sl, payload, x, y, w)
        elif typ == "code":
            # framed terminal panel (white fill, thin accent border) --
            # matches Dave's hand-framed slide-9 code box, pass 2
            maxlen = max(len(l) for l in payload) or 1
            size = min(11.5, max(8, w * 72.0 / (0.62 * maxlen)))
            hgt = len(payload) * size * 1.32 / 72.0 + 0.18
            sh = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x),
                                     Inches(y), Inches(w), Inches(hgt))
            sh.fill.solid(); sh.fill.fore_color.rgb = WHITE
            sh.line.color.rgb = ACCENT; sh.line.width = Pt(0.75)
            sh.shadow.inherit = False
            tf = sh.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.TOP
            tf.margin_left = tf.margin_right = Inches(0.10)
            tf.margin_top = tf.margin_bottom = Inches(0.05)
            for j, l in enumerate(payload):
                para(tf, l if l else " ", size=size, mono=True,
                     first=(j == 0), space_after=1)
            y += hgt + 0.12
        elif typ == "img":
            cap, path, h = payload
            exact = bool(ov and "w" in ov and "h" in ov)
            avail = ov["h"] if exact else \
                (min(h, BOTTOM - y) if h else (BOTTOM - y))
            y = place_img(sl, path, x, y, w, avail, cap, exact=exact,
                          cap_ov=ov)
        elif typ == "stack":
            n = len(payload)
            bh, gap = 0.78, 0.14
            for j, (hd, sub) in enumerate(payload):
                sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(x), Inches(y),
                                         Inches(w), Inches(bh))
                sh.adjustments[0] = 0.10
                sh.fill.solid(); sh.fill.fore_color.rgb = LIGHT
                sh.line.color.rgb = ACCENT; sh.line.width = Pt(1.0)
                sh.shadow.inherit = False
                tf = sh.text_frame; tf.word_wrap = True
                tf.margin_left = tf.margin_right = Inches(0.08)
                tf.margin_top = tf.margin_bottom = Inches(0.03)
                para(tf, hd, size=11.5, bold=True, first=True,
                     align=PP_ALIGN.CENTER, space_after=1)
                if sub:
                    para(tf, sub, size=10, color=INK,
                         align=PP_ALIGN.CENTER, space_after=0)
                if j < n - 1:
                    ar = sl.shapes.add_shape(
                        MSO_SHAPE.DOWN_ARROW,
                        Inches(x + w / 2 - 0.09), Inches(y + bh - 0.02),
                        Inches(0.18), Inches(gap + 0.04))
                    ar.fill.solid(); ar.fill.fore_color.rgb = ACCENT
                    ar.line.fill.background(); ar.shadow.inherit = False
                y += bh + gap
            y += 0.05
        elif typ == "cards":
            ncol, items = payload
            gap_x, gap_y, ch = 0.28, 0.22, 1.35
            cw = (w - (ncol - 1) * gap_x) / ncol
            for j, card in enumerate(items):
                r, c = divmod(j, ncol)
                cx = x + c * (cw + gap_x)
                cy = y + r * (ch + gap_y)
                sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(cx), Inches(cy),
                                         Inches(cw), Inches(ch))
                sh.adjustments[0] = 0.07
                sh.fill.solid(); sh.fill.fore_color.rgb = LIGHT
                sh.line.color.rgb = ACCENT; sh.line.width = Pt(1.0)
                sh.shadow.inherit = False
                tf = sh.text_frame; tf.word_wrap = True
                tf.margin_left = tf.margin_right = Inches(0.08)
                tf.margin_top = tf.margin_bottom = Inches(0.04)
                for k, linetxt in enumerate(card):
                    para(tf, linetxt, size=10.5, bold=(k == 0),
                         first=(k == 0), align=PP_ALIGN.CENTER,
                         space_after=1)
            rows = ceil(len(items) / ncol)
            y += rows * ch + (rows - 1) * gap_y + 0.12
        cur[colname] = y
        if colname == "full":
            for c in cur:
                cur[c] = y


def render_title(t):
    sl = prs.slides.add_slide(BLANK)
    tf = tb(sl, 0.9, 1.9, SW_IN - 1.8, 3.2)
    para(tf, t["title"], size=40, bold=True, color=ACCENT, first=True)
    for b in t["big"]:
        para(tf, b, size=22, color=INK, space_after=14)
    for p in t["paras"]:
        para(tf, p, size=15, color=GRAY, space_after=8)
    tf2 = tb(sl, 0.9, 6.0, SW_IN - 1.8, 1.0)
    for j, n in enumerate(t["notes"]):
        para(tf2, n, size=13, color=GRAY, first=(j == 0), space_after=4)


def main():
    with open(MD, encoding="utf-8") as f:
        title_slide, slides = parse(f.read())
    render_title(title_slide)
    for spec in slides:
        render_slide(spec)
    prs.save(OUT)
    print(f"wrote {OUT}  ({len(slides) + 1} slides)")


if __name__ == "__main__":
    sys.exit(main())
